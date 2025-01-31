# Recriar e salvar a apresentação novamente

from pptx import Presentation
from pptx.util import Inches

# Criar uma nova apresentação
prs = Presentation()

# ---- SLIDE 1: TÍTULO ----
slide_layout = prs.slide_layouts[0]  # Layout de título
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "📊 Estratégia de Precificação para Aluguéis Temporários em Nova York"
subtitle.text = "Análise Exploratória e Modelagem Preditiva de Preços\nIndicium - Cliente Confidencial"

# ---- SLIDE 2: CONTEXTO ----
slide_layout = prs.slide_layouts[1]  # Layout de título e conteúdo
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "📍 Contexto do Projeto"
content.text = (
    "🔹 A Indicium foi contratada para apoiar a criação de uma plataforma de aluguéis temporários em Nova York.\n"
    "🔹 O objetivo é entender a precificação do mercado, explorando dados de um concorrente.\n"
    "🔹 Além disso, buscamos desenvolver um modelo preditivo para estimar preços de aluguéis."
)

# ---- SLIDE 3: EXPLORAÇÃO DOS DADOS (EDA) ----
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "🔍 Exploração dos Dados (EDA)"
content.text = (
    "📌 Principais descobertas:\n"
    "✅ Manhattan é o bairro mais caro, seguido por Brooklyn.\n"
    "✅ Apartamentos inteiros têm preços significativamente mais altos.\n"
    "✅ Propriedades com menos disponibilidade no ano tendem a cobrar mais.\n"
    "✅ Palavras como 'Luxury', 'View' e 'Central' aparecem nos imóveis mais caros."
)

# ---- SLIDE 4: ONDE INVESTIR? ----
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "🏡 Onde Investir?"
content.text = (
    "🔹 **Bairro:** Manhattan é a melhor escolha para maximizar ganhos.\n"
    "🔹 **Tipo de imóvel:** Alugar um apartamento inteiro é mais rentável.\n"
    "🔹 **Descrição do anúncio:** Destacar localização e exclusividade aumenta o valor percebido.\n"
    "🔹 **Disponibilidade:** Limitar os dias disponíveis pode justificar preços mais altos."
)

# ---- SLIDE 5: IMPACTO DE NOITES MÍNIMAS E DISPONIBILIDADE ----
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "⏳ O Número de Noites Mínimas e Disponibilidade Importam?"
content.text = (
    "📌 **Conclusões:**\n"
    "✅ Noites mínimas não têm grande impacto no preço.\n"
    "✅ Propriedades disponíveis o ano todo tendem a cobrar menos.\n"
    "✅ Menos disponibilidade pode indicar exclusividade e permitir preços mais altos."
)

# ---- SLIDE 6: PADRÃO NO NOME DOS IMÓVEIS ----
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "📢 Existe um Padrão no Nome dos Imóveis Mais Caros?"
content.text = (
    "🔍 Sim! Palavras-chave que aparecem frequentemente em imóveis mais caros:\n"
    "✅ 'Luxury'\n"
    "✅ 'View'\n"
    "✅ 'Central'\n"
    "✅ 'Exclusive'\n"
    "💡 **Sugestão:** Usar essas palavras na descrição pode aumentar o valor percebido."
)

# ---- SLIDE 7: MODELAGEM PREDITIVA ----
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "🤖 Modelagem Preditiva: Como Prevemos os Preços?"
content.text = (
    "📌 Utilizamos um modelo de Machine Learning para prever preços de aluguéis.\n"
    "🔹 **Tipo de problema:** Regressão (valor numérico do preço).\n"
    "🔹 **Modelo utilizado:** Random Forest Regressor.\n"
    "🔹 **Principais variáveis:** Localização, tipo de imóvel, número de reviews, disponibilidade, entre outras."
)

# ---- SLIDE 8: RESULTADOS DO MODELO ----
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "📈 Resultados do Modelo"
content.text = (
    "📌 **Métricas de desempenho:**\n"
    "✅ **Erro Quadrático Médio (MSE):** 6.936,24\n"
    "✅ **R² Score:** 0,45 (o modelo explica cerca de 45% da variação dos preços).\n"
    "💡 **Possíveis melhorias:** Mais dados, refinamento de variáveis e novos algoritmos."
)

# ---- SLIDE 9: PREVISÃO DE PREÇO PARA UM IMÓVEL ----
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "💰 Previsão de Preço para um Imóvel"
content.text = (
    "📍 **Características do imóvel:**\n"
    "✅ Localizado em Midtown Manhattan.\n"
    "✅ Apartamento inteiro, 45 reviews, 355 dias disponíveis no ano.\n"
    "📌 **Preço previsto:** **$250 por noite**\n"
    "💡 Pode valer a pena testar um preço mais alto e avaliar a demanda."
)

# ---- SLIDE 10: CONCLUSÃO ----
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "🏁 Conclusão e Próximos Passos"
content.text = (
    "✅ **Onde investir?** Manhattan, com apartamentos inteiros.\n"
    "✅ **Como se destacar?** Descrição atrativa e disponibilidade estratégica.\n"
    "✅ **Modelo preditivo:** Já oferece boas estimativas, mas pode ser aprimorado.\n"
    "🚀 **Próximos passos:** Refinar o modelo, coletar mais dados e testar estratégias de precificação."
)

# Salvar a apresentação
pptx_path = "/mnt/data/Analise_Precificacao_NY.pptx"
prs.save(r"C:\Projetos\DesafioBruna\Analise_Precificacao_NY.pptx")
